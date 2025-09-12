from pptx import Presentation

# Create a presentation object
prs = Presentation()

# Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Portfolio - HEMA PRIYA M"
subtitle.text = "Aspiring Data Analyst | B.Sc Computer Science"

# About Me
about_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(about_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "About Me"
content.text = (
    "I am an enthusiastic and self-motivated B.Sc Computer Science candidate with analytical "
    "and problem-solving skills. I possess basic knowledge in MS Excel and am able to work both "
    "in a team environment as well as independently."
)

# Skills
skills_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(skills_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Skills"
content.text = (
    "- MS Office\n"
    "- Designing\n"
    "- Data Analytics (Basics)\n"
    "- Communication Skills\n"
    "- Languages: English, Tamil"
)

# Education
edu_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(edu_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Education"
content.text = (
    "- SSLC: Sharon Matric Higher Secondary School - 75.2%\n"
    "- HSC: Sethu Bhaskara Matric Higher Secondary School - 68.8% (2023-2024)\n"
    "- B.Sc Computer Science (Pursuing) - S.A Arts Science and College, Madras University"
)

# Project
proj_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(proj_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Project 1: Student Record Management System"
content.text = (
    "- Developed a simple application to store and manage student details.\n"
    "- Used MS Excel and Python for data handling.\n"
    "- Designed with a focus on user-friendly interface and accuracy.\n"
    "- Learned teamwork, documentation, and problem-solving."
)

# Certificates
cert_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(cert_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Certificates"
content.text = (
    "- Diploma in Computer Application\n"
    "- Masterclass in Data Analytics"
)

# Contact
contact_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(contact_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Contact"
content.text = (
    "Phone: 9884233943\n"
    "Email: hemamoorthy8586@gmail.com\n"
    "Location: Chennai, India"
)

# Save the presentation
prs.save("Portfolio_HemaPriyaM.pptx")
